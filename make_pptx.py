#!/usr/bin/env python3
"""Generate 31-slide EIS project PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSAT, MSO_CONNECTOR_TYPE

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W, H = 13.33, 7.5
BLANK = prs.slide_layouts[6]

NAVY   = RGBColor(0x0F, 0x1F, 0x3D)
BLUE   = RGBColor(0x15, 0x65, 0xC0)
TEAL   = RGBColor(0x00, 0x89, 0x7B)
CYAN   = RGBColor(0x00, 0xBC, 0xD4)
AMBER  = RGBColor(0xFF, 0x8F, 0x00)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xEF, 0x6C, 0x00)
RED    = RGBColor(0xC6, 0x28, 0x28)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
BROWN  = RGBColor(0x4E, 0x34, 0x2E)
SLATE  = RGBColor(0x37, 0x47, 0x4F)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF4, 0xF7, 0xFB)
BODY   = RGBColor(0x37, 0x41, 0x51)
MUTED  = RGBColor(0x64, 0x74, 0x8B)
CARD   = RGBColor(0xE8, 0xF0, 0xFE)
LCARD  = RGBColor(0xF0, 0xF5, 0xFF)

def bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def _sh(slide, kind, x, y, w, h, fill, border=None, bw=0.75):
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(bw)
    else: s.line.fill.background()
    return s

def box(slide, x, y, w, h, fill, border=None, bw=0.75):
    return _sh(slide, MSAT.RECTANGLE, x, y, w, h, fill, border, bw)

def rbox(slide, x, y, w, h, fill, border=None, bw=0.75):
    return _sh(slide, MSAT.ROUNDED_RECTANGLE, x, y, w, h, fill, border, bw)

def oval(slide, x, y, w, h, fill):
    return _sh(slide, MSAT.OVAL, x, y, w, h, fill)

def lbox(slide, text, x, y, w, h, fill, tc=WHITE, fs=10, bold=True,
         border=None, bw=0.75, align=PP_ALIGN.CENTER, italic=False, rounded=True):
    s = rbox(slide,x,y,w,h,fill,border,bw) if rounded else box(slide,x,y,w,h,fill,border,bw)
    tf=s.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(fs); r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=tc
    return s

def lbox2(slide, title, subtitle, x, y, w, h, fill, tc=WHITE, tfs=11, sfs=9, border=None, bw=0.75):
    s = rbox(slide,x,y,w,h,fill,border,bw)
    tf=s.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=title; r.font.size=Pt(tfs); r.font.bold=True; r.font.color.rgb=tc
    p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER
    r2=p2.add_run(); r2.text=subtitle; r2.font.size=Pt(sfs); r2.font.bold=False; r2.font.color.rgb=tc
    return s

def t(slide, text, x, y, w, h, fs=12, c=BODY, bold=False,
      align=PP_ALIGN.LEFT, italic=False, wrap=True):
    b=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=b.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(fs); r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=c
    return b

def txt(slide, lines, x, y, w, h, wrap=True):
    b=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=b.text_frame; tf.word_wrap=wrap; first=True
    for (text,fs,color,bold,italic,align) in lines:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align if align else PP_ALIGN.LEFT
        r=p.add_run(); r.text=text
        r.font.size=Pt(fs); r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color
    return b

def cn(slide, x1, y1, x2, y2, color=MUTED, w=1.5):
    c=slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
      Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    c.line.color.rgb=color; c.line.width=Pt(w); return c

def hdr(slide, title, sub=None):
    t(slide,title,0.4,0.18,W-0.8,0.7,fs=26,c=NAVY,bold=True)
    box(slide,0.4,0.88,W-0.8,0.04,BLUE)
    if sub: t(slide,sub,0.4,0.95,W-0.8,0.38,fs=11,c=MUTED,italic=True)

def dark_hdr(slide, title, sub=None):
    """Header for navy-bg slides."""
    t(slide,title,0.6,0.2,W-1.2,0.75,fs=28,c=WHITE,bold=True)
    box(slide,0.6,0.94,W-1.2,0.04,TEAL)
    if sub: t(slide,sub,0.6,1.0,W-1.2,0.38,fs=12,c=MUTED,italic=True)

# ═══ SLIDE 1 — Title ══════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,NAVY)
oval(sl,-1.8,-1.8,5.5,5.5,TEAL); oval(sl,10.5,5.2,4.0,4.0,BLUE)
t(sl,"Electrochemical Impedance Spectroscopy",1.0,1.5,11.33,1.1,fs=38,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
t(sl,"EIS Measurement System",1.0,2.65,11.33,0.85,fs=30,c=CYAN,align=PP_ALIGN.CENTER)
box(sl,4.3,3.6,4.73,0.04,TEAL)
t(sl,"Master's Thesis  |  STM32F303CCT6  |  Python Desktop GUI",1.0,3.75,11.33,0.5,fs=14,c=MUTED,italic=True,align=PP_ALIGN.CENTER)
t(sl,"Portable swept-sine EIS instrument for electrochemical electrode characterization",1.0,4.28,11.33,0.5,fs=12,c=RGBColor(0x8B,0x9B,0xB4),align=PP_ALIGN.CENTER)
t(sl,"2026",0.4,6.95,W-0.8,0.35,fs=11,c=MUTED,align=PP_ALIGN.RIGHT)

# ═══ SLIDE 2 — Project Overview ═══════════════════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Project Overview","A compact USB-controlled EIS instrument for master's thesis electrochemical research")
desc=[
    ("What it is",13,NAVY,True,False,None),
    ("A portable electrochemical impedance spectroscopy (EIS) instrument built around the STM32F303CCT6 microcontroller. "
     "It sweeps a sinusoidal excitation signal across a user-defined frequency range and measures complex impedance at each step.",11,BODY,False,False,None),
    ("",5,BODY,False,False,None),
    ("How it works",13,NAVY,True,False,None),
    ("An AD9833 DDS chip generates the stimulus sine wave. The signal passes through a TIA stage and both the reference and "
     "TIA channels are simultaneously digitized by a dual-channel ADC with DMA (256 samples each). A least-squares tone "
     "estimator extracts magnitude and phase at the exact measurement frequency.",11,BODY,False,False,None),
    ("",5,BODY,False,False,None),
    ("Key outputs",13,NAVY,True,False,None),
    ("Bode plots (|Z| and phase vs. frequency) and a Nyquist plot (Re(Z) vs. -Im(Z)) are displayed live in the Python "
     "desktop GUI. All sweep data is also logged to on-board W25Q32 SPI flash and can be exported as CSV.",11,BODY,False,False,None),
]
txt(sl,desc,0.4,1.4,6.7,5.8)
pipeline=[("PC — Python GUI (PyQt5)",BLUE),("USB-UART  (CP2102 Bridge)",TEAL),
          ("STM32F303CCT6 Firmware",NAVY),("AD9833 DDS  ->  Signal Chain",ORANGE),
          ("ADC + Tone Estimator  ->  Z(f)",GREEN)]
bx_x,bw2,bh2,bgap=7.6,5.3,0.68,0.22; by=1.4
for i,(label,fill) in enumerate(pipeline):
    lbox(sl,label,bx_x,by,bw2,bh2,fill,tc=WHITE,fs=12,bold=True)
    if i<len(pipeline)-1: cn(sl,bx_x+bw2/2,by+bh2,bx_x+bw2/2,by+bh2+bgap+0.01,color=BLUE,w=2)
    by+=bh2+bgap+0.28

# ═══ SLIDE 3 (NEW) — EIS Theory & Measurement Principle ═══════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"EIS Theory & Measurement Principle","Electrochemical Impedance Spectroscopy — the physics behind the measurement")
txt(sl,[
    ("What is Impedance?",13,NAVY,True,False,None),
    ("Impedance Z(f) = V(f) / I(f) — the complex ratio of voltage to current at each frequency. "
     "It extends resistance to the complex plane: Z = Re(Z) + j·Im(Z).",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Swept-Sine Method",13,NAVY,True,False,None),
    ("A pure sine wave at frequency f is applied to the electrochemical cell. The resulting current "
     "(or voltage across a reference resistor) is measured. Magnitude |Z| = V_exc / I_resp and "
     "phase angle phi = angle(V) - angle(I) are extracted at each frequency step.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Randles Circuit Model",13,NAVY,True,False,None),
    ("A simple electrochemical cell can be modelled as: solution resistance Rs in series with "
     "the parallel combination of double-layer capacitance Cdl and charge-transfer resistance Rct. "
     "This gives a characteristic semicircle on the Nyquist plot.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Applications",13,NAVY,True,False,None),
    ("Battery state-of-health, corrosion monitoring, biosensor characterization, "
     "membrane impedance, supercapacitor analysis.",11,BODY,False,False,None),
],0.4,1.38,6.0,5.9)
# Right: diagram of Randles circuit + Nyquist sketch
rbox(sl,6.7,1.38,6.2,2.8,WHITE,border=BLUE,bw=1)
t(sl,"Randles Equivalent Circuit",6.9,1.45,5.8,0.35,fs=12,c=NAVY,bold=True,align=PP_ALIGN.CENTER)
# Rs box
lbox(sl,"Rs\n(Solution\nResistance)",7.1,1.95,1.5,1.2,BLUE,tc=WHITE,fs=10)
# Cdl box
lbox(sl,"Cdl\n(Double-layer\nCapacitance)",9.0,1.82,1.5,0.9,TEAL,tc=WHITE,fs=9)
# Rct box
lbox(sl,"Rct\n(Charge-transfer\nResistance)",9.0,2.8,1.5,0.9,ORANGE,tc=WHITE,fs=9)
# Connect lines
cn(sl,8.6,2.55,9.0,2.27,color=MUTED,w=1)
cn(sl,8.6,2.55,9.0,3.25,color=MUTED,w=1)
cn(sl,10.5,2.27,10.9,2.55,color=MUTED,w=1)
cn(sl,10.5,3.25,10.9,2.55,color=MUTED,w=1)
lbox(sl,"Electrode / DUT",11.1,1.95,1.5,1.2,RGBColor(0x4E,0x34,0x2E),tc=WHITE,fs=9)
t(sl,"Current I",6.95,2.55,1.0,0.28,fs=8,c=MUTED,italic=True)
# Nyquist sketch
rbox(sl,6.7,4.35,6.2,2.8,WHITE,border=TEAL,bw=1)
t(sl,"Nyquist Plot — Randles Semicircle",6.9,4.42,5.8,0.35,fs=12,c=NAVY,bold=True,align=PP_ALIGN.CENTER)
t(sl,"Re(Z) ->",7.0,6.8,2.0,0.3,fs=9,c=MUTED,italic=True)
t(sl,"^ -Im(Z)",6.7,4.8,1.0,0.3,fs=9,c=MUTED,italic=True)
# Draw semicircle hint with boxes
lbox(sl,"Rs",7.3,6.35,0.6,0.5,BLUE,tc=WHITE,fs=8,bold=False)
lbox(sl,"Rs+Rct",11.7,6.35,1.1,0.5,ORANGE,tc=WHITE,fs=8,bold=False)
t(sl,"<--- Rct diameter --->",8.2,6.4,3.2,0.35,fs=9,c=TEAL,align=PP_ALIGN.CENTER)
cn(sl,7.6,6.35,9.3,5.0,color=TEAL,w=1.5)
cn(sl,9.3,5.0,11.7,6.35,color=TEAL,w=1.5)

# ═══ SLIDE 4 (NEW) — System Specifications ════════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"System Specifications & Key Parameters","Hardware and software performance figures for the EIS instrument")
specs=[
    ("Frequency range","1 Hz – 100 kHz (software limit; ADC Nyquist ~121 kHz)",BLUE),
    ("ADC sample rate","72 MHz / ((61.5+12.5)x2) = 486 kHz total; 243 kHz per channel",NAVY),
    ("Samples per step","256 per channel (512 ADC words per frequency point via DMA)",TEAL),
    ("DDS resolution","25 MHz / 2^28 = 0.093 Hz  (AD9833 with 25 MHz MCLK)",GREEN),
    ("Impedance formula","Z = (|V_ref| / |V_sig|) x Rf  ;  phi = angle(V_sig) - angle(V_ref)",ORANGE),
    ("Flash storage","~390 records in 16 KB log region (4 x 4 KB sectors of W25Q32)",PURPLE),
    ("UART protocol","115 200 baud, 8N1, ASCII line protocol (CR+LF terminated)",RED),
    ("Supply voltage","USB-C 5 V -> MCP1700 LDO -> 3.3 V  (250 mA max)",BROWN),
    ("Reference voltage","ADR4533BRZ 3.3 V precision ref for STM32 VDDA and VREF+",SLATE),
    ("System clock","8 MHz HSE crystal -> STM32 PLL x9 = 72 MHz",RGBColor(0x2E,0x7D,0x32)),
    ("SPI clock","36 MHz (SPI1 prescaler /2 from 72 MHz for AD9833 and W25Q32)",RGBColor(0x1A,0x23,0x7E)),
    ("Phase accuracy","Determined by ADC sample timing precision and 486 kHz sample rate",AMBER),
]
cols=2; rw=(W-0.9)/cols; rh=0.55
for i,(param,value,fill) in enumerate(specs):
    col=i%cols; row=i//cols
    sx=0.4+col*(rw+0.1); sy=1.4+row*(rh+0.08)
    lbox2(sl,param,value,sx,sy,rw,rh,fill,tc=WHITE,tfs=11,sfs=9)

# ═══ SLIDE 5 — Python Desktop GUI ════════════════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Python Desktop GUI","PyQt5  |  Matplotlib  |  PySerial  |  NumPy — cross-platform, install via pip")
txt(sl,[
    ("Connection & Control",12,NAVY,True,False,None),
    ("Port selector, baud, Connect/Disconnect. Start Sweep, Stop, Status, Flash Status, Dump Flash, Erase Flash.",10,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Sweep Configuration",12,NAVY,True,False,None),
    ("Start freq, stop freq, step freq, feedback Rf — applied live via SET_* UART commands.",10,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Live Plots",12,NAVY,True,False,None),
    ("Bode |Z| (log-log), Bode phase (log-x), Nyquist (Re vs -Im) — all update point-by-point during sweep.",10,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Status Panel",12,NAVY,True,False,None),
    ("Live state (IDLE/RUNNING), current frequency, points done, flash record count & capacity.",10,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Color-Coded Serial Console",12,NAVY,True,False,None),
    ("Green = received, Blue = sent, Red = errors. Full UART log visible at all times.",10,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Data Export",12,NAVY,True,False,None),
    ("One-click CSV export: freq, |Z|, phase, Re(Z), Im(Z). Flash dump also populates plots.",10,BODY,False,False,None),
],0.4,1.38,5.5,5.9)
mx,mw,my,mh=6.2,6.8,1.35,5.9
box(sl,mx,my,mw,mh,RGBColor(0xDC,0xDC,0xDC),border=MUTED,bw=1)
box(sl,mx,my,mw,0.48,NAVY)
t(sl,"Port: COM5  [Connect]  [Start]  [Stop]  [Status]  [Flash]  [Export]",mx+0.08,my+0.12,mw-0.16,0.28,fs=7,c=WHITE)
sbw=1.65
box(sl,mx,my+0.48,sbw,mh-0.48,CARD)
t(sl,"Sweep Config",mx+0.06,my+0.52,sbw-0.1,0.22,fs=7.5,c=NAVY,bold=True)
t(sl,"Start: 1000 Hz\nStop: 100000 Hz\nStep: 1000 Hz\nRf: 10000\n[Apply]",mx+0.08,my+0.76,sbw-0.14,0.85,fs=7,c=BODY)
t(sl,"Live Status",mx+0.06,my+1.66,sbw-0.1,0.22,fs=7.5,c=NAVY,bold=True)
t(sl,"State: RUNNING\nFreq: 5000 Hz\nPts: 4/99\nFlash: 4 recs",mx+0.08,my+1.90,sbw-0.14,0.75,fs=7,c=BODY)
t(sl,"Last Point",mx+0.06,my+2.70,sbw-0.1,0.22,fs=7.5,c=NAVY,bold=True)
t(sl,"|Z|: 1234 O\nPhase: -23.5\nRe: 1134 O\nIm: -497 O",mx+0.08,my+2.95,sbw-0.14,0.70,fs=7,c=BODY)
px=mx+sbw; pw=mw-sbw; ph_top=(mh-0.48)*0.56
box(sl,px,my+0.48,pw/2,ph_top,WHITE,border=MUTED,bw=0.4)
t(sl,"Bode |Z|",px+0.06,my+0.54,pw/2-0.1,0.2,fs=7.5,c=NAVY,bold=True)
cn(sl,px+0.18,my+0.48+ph_top*0.85,px+pw/2-0.18,my+0.48+ph_top*0.2,color=BLUE,w=1.2)
box(sl,px+pw/2,my+0.48,pw/2,ph_top,WHITE,border=MUTED,bw=0.4)
t(sl,"Bode Phase",px+pw/2+0.06,my+0.54,pw/2-0.1,0.2,fs=7.5,c=NAVY,bold=True)
cn(sl,px+pw/2+0.18,my+0.48+ph_top*0.2,px+pw-0.18,my+0.48+ph_top*0.85,color=TEAL,w=1.2)
nyq_y=my+0.48+ph_top; nyq_h=(mh-0.48)*0.32
box(sl,px,nyq_y,pw,nyq_h,WHITE,border=MUTED,bw=0.4)
t(sl,"Nyquist",px+0.06,nyq_y+0.05,pw-0.1,0.2,fs=7.5,c=NAVY,bold=True)
con_y=my+mh-0.58; box(sl,mx,con_y,mw,0.58,RGBColor(0x1E,0x1E,0x1E))
t(sl,"<- DATA,5000,1234.56,-23.45    -> START    <- SWEEP,BEGIN",mx+0.06,con_y+0.12,mw-0.12,0.35,fs=6.5,c=RGBColor(0x4C,0xAF,0x50))

# ═══ SLIDE 6 (NEW) — GUI Architecture & Module Connections ════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"GUI Architecture & Module Connections","Five Python modules — how they relate and communicate")
# Module descriptions left
txt(sl,[
    ("main.py — Application Entry Point",12,NAVY,True,False,None),
    ("QMainWindow that owns the toolbar, central layout, and status bar. "
     "It instantiates all other modules and wires their Qt signals to slots. "
     "Parses every incoming UART line and routes it to DataModel, PlotWidget, and ConfigPanel.",10,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("serial_worker.py — Background I/O",12,NAVY,True,False,None),
    ("A QThread that opens the COM port, reads bytes continuously, assembles "
     "line-by-line output, and emits line_received(str). Also handles send_command(). "
     "Runs entirely off the GUI thread so the UI never freezes.",10,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("data_model.py — Data Store",12,NAVY,True,False,None),
    ("Holds a Python list of ImpedancePoint dataclasses. Provides add_point(), "
     "clear(), get_arrays(), count(), and export_csv(). No Qt dependency.",10,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("plot_widget.py — Live Visualization",12,NAVY,True,False,None),
    ("Three Matplotlib FigureCanvas objects in a QSplitter. update_plots() redraws "
     "all three (Bode |Z|, Bode phase, Nyquist) on every new DATA point.",10,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("config_panel.py — Sidebar",12,NAVY,True,False,None),
    ("QGroupBox sections for sweep parameters, live status labels, and last-point readout. "
     "Emits settings_applied(dict) when the user clicks Apply.",10,BODY,False,False,None),
],0.4,1.38,6.0,5.9)
# Right: module dependency diagram
mods6={
    "main.py":          (9.5,2.0,NAVY),
    "serial_worker.py": (7.2,3.3,TEAL),
    "data_model.py":    (9.5,3.3,BLUE),
    "plot_widget.py":   (11.8,3.3,GREEN),
    "config_panel.py":  (9.5,4.6,ORANGE),
    "COM Port\n(PySerial)":(7.2,4.6,SLATE),
}
mw6,mh6=1.9,0.65
for name,(mx6,my6,fill) in mods6.items():
    lbox(sl,name,mx6,my6,mw6,mh6,fill,tc=WHITE,fs=9,bold=True)
# Connections
def mc6(n): mx6,my6,_=mods6[n]; return mx6+mw6/2,my6+mh6/2
edges6=[("main.py","serial_worker.py"),("main.py","data_model.py"),
        ("main.py","plot_widget.py"),("main.py","config_panel.py"),
        ("serial_worker.py","COM Port\n(PySerial)")]
for a,b in edges6:
    ax,ay=mc6(a); bx,by=mc6(b)
    cn(sl,ax,ay,bx,by,color=MUTED,w=1.2)
t(sl,"line_received signal",7.6,2.8,2.5,0.32,fs=8,c=TEAL,italic=True)
t(sl,"settings_applied signal",8.5,4.28,2.5,0.32,fs=8,c=ORANGE,italic=True)
t(sl,"update_plots() call",10.7,2.65,2.0,0.32,fs=8,c=GREEN,italic=True)

# ═══ SLIDE 7 (NEW) — GUI User Workflow ════════════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"GUI User Workflow — Step by Step","Complete interaction sequence from power-on to data export")
steps7=[
    ("1. Launch",   "Run:  python gui/main.py\nPyQt5 window opens. Serial console is blank.",           NAVY),
    ("2. Select Port","Click the port dropdown — available COM ports are listed. Pick the CP2102 port (e.g. COM5).",BLUE),
    ("3. Connect",  "Click Connect. SerialWorker opens the port at 115200 baud. Status bar shows Connected.",     TEAL),
    ("4. Configure","Enter start/stop/step frequency and Rf in the sidebar. Click Apply — SET_* commands are sent.", GREEN),
    ("5. Start Sweep","Click Start Sweep. GUI clears data, sends START. STM32 replies SWEEP,BEGIN.",              ORANGE),
    ("6. Watch Plots","DATA lines arrive — Bode and Nyquist plots build up point by point in real time.",         PURPLE),
    ("7. Stop / Wait","Click Stop (sends STOP) or wait for SWEEP,DONE. Status changes back to IDLE.",            RED),
    ("8. Export",   "Click Export CSV — file dialog opens. CSV is written with freq, |Z|, phase, Re(Z), Im(Z).", BROWN),
    ("9. Flash Ops","Click Flash Status to see stored records. Dump Flash to replay a previous sweep.",           SLATE),
]
cols7=3; rw7=(W-0.9)/cols7; rh7=1.55
for i,(step,desc,fill) in enumerate(steps7):
    col=i%cols7; row=i//cols7
    sx=0.4+col*(rw7+0.15); sy=1.4+row*(rh7+0.12)
    lbox2(sl,step,desc,sx,sy,rw7,rh7,fill,tc=WHITE,tfs=12,sfs=9.5)

# ═══ SLIDE 8 — Hardware Components ═══════════════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Hardware Components","Key ICs and components on the EIS PCB")
components=[
    ("STM32F303CCT6","72 MHz ARM Cortex-M4, dual ADC with DMA, SPI, UART — central controller",NAVY),
    ("AD9833 DDS","Programmable sine waveform generator; 25 MHz MCLK; SPI-controlled to 12.5 MHz",BLUE),
    ("OPA2140AID (x2)","Dual precision FET-input op-amp: Amp1A = excitation buffer, Amp2A = TIA",TEAL),
    ("ADR4533BRZ","Precision 3.3 V voltage reference (<25 uV/degC) feeds STM32 VDDA and VREF+",ORANGE),
    ("W25Q32JVS","32 Mbit SPI flash; shares SPI1 via separate CS on PB12; stores sweep logs",GREEN),
    ("CP2102","USB-UART bridge; STM32 USART1 (PA9 TX / PA10 RX) at 115 200 baud",RED),
    ("MCP1700 LDO","Ultra-low IQ LDO: USB-C VBUS 5 V -> 3.3 V system supply, 250 mA",PURPLE),
    ("ASE-25.000 MHz","Crystal oscillator for AD9833 MCLK — DDS frequency resolution 0.093 Hz",RGBColor(0x00,0x69,0x8D)),
    ("ABM3B-8.000 MHz","Crystal for STM32 HSE input; on-chip PLL x9 = 72 MHz system clock",BROWN),
]
cols8=3; cw8=(W-0.8)/cols8; ch8=(H-1.5)/cols8
for i,(name,desc,fill) in enumerate(components):
    col=i%cols8; row=i//cols8
    lbox2(sl,name,desc,0.4+col*cw8+0.1,1.5+row*ch8+0.1,cw8-0.2,ch8-0.2,fill,tc=WHITE,tfs=12,sfs=9.5)

# ═══ SLIDE 9 (NEW) — STM32 & AD9833 In Detail ════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"STM32F303CCT6 & AD9833 — In Detail","The two core digital ICs: microcontroller and DDS signal generator")
txt(sl,[
    ("STM32F303CCT6 — Microcontroller",13,NAVY,True,False,None),
    ("Core: 72 MHz ARM Cortex-M4F with FPU and DSP instructions.",11,BODY,False,False,None),
    ("Clock: 8 MHz HSE crystal -> PLL multiplier x9 = 72 MHz system clock.",11,BODY,False,False,None),
    ("ADC1: 12-bit, dual-channel scan mode (PA0 + PA1), triggered by software, "
     "transferred via DMA1 Channel 1. Sample time = 61.5 cycles, conversion = 12.5 cycles.",11,BODY,False,False,None),
    ("SPI1: PA5=SCK, PA6=MISO, PA7=MOSI. Two chip-selects: PA4 (AD9833 FSYNC), PB12 (W25Q32 CS).",11,BODY,False,False,None),
    ("USART1: PA9=TX, PA10=RX at 115200 baud 8N1. DMA or interrupt driven receive.",11,BODY,False,False,None),
    ("GPIO: PB13=STATUS LED, PB12=W25Q32 CS, PA4=AD9833 FSYNC (soft-controlled).",11,BODY,False,False,None),
    ("",5,BODY,False,False,None),
    ("AD9833 — DDS Signal Generator",13,NAVY,True,False,None),
    ("25 MHz MCLK from ASE-25.000MHz crystal oscillator.",11,BODY,False,False,None),
    ("Frequency word = round(f_out x 2^28 / 25e6), 28-bit value split into two "
     "14-bit SPI writes (D15=1, D14=0 for FREQ0 register, first write sets LSB, second MSB).",11,BODY,False,False,None),
    ("Output: sine wave approx 0.6 Vpp. FSYNC line low during SPI transfer, high to latch.",11,BODY,False,False,None),
    ("Minimum frequency step: 25e6 / 2^28 = 0.093 Hz.",11,BODY,False,False,None),
],0.4,1.38,6.5,5.9)
# Right: pin diagram blocks
lbox(sl,"STM32F303CCT6",7.3,1.5,5.6,1.0,NAVY,tc=WHITE,fs=13,bold=True)
pins=[("PA0 — ADC1 CH1 (Ref)",BLUE),("PA1 — ADC1 CH2 (TIA)",BLUE),
      ("PA4 — AD9833 FSYNC",TEAL),("PA5/6/7 — SPI1 CLK/MISO/MOSI",TEAL),
      ("PA9/PA10 — USART1 TX/RX",GREEN),("PB12 — W25Q32 CS",GREEN),
      ("PB13 — Status LED",ORANGE)]
py9=2.6
for label,fill in pins:
    lbox(sl,label,7.3,py9,5.6,0.46,fill,tc=WHITE,fs=10,bold=False,align=PP_ALIGN.LEFT); py9+=0.48

# ═══ SLIDE 10 (NEW) — Analog Front-End & Power ════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Analog Front-End & Power Management","OPA2140 op-amps, ADR4533 reference, and MCP1700 LDO regulator")
txt(sl,[
    ("OPA2140AID — Excitation Buffer (Amp1A)",13,NAVY,True,False,None),
    ("Unity-gain voltage follower. High input impedance (1 TOhm) prevents loading the AD9833 output. "
     "Low output impedance drives the electrode with minimal voltage drop.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("OPA2140AID — Transimpedance Amplifier (Amp2A)",13,NAVY,True,False,None),
    ("Converts current through the DUT to voltage: V_out = -I_DUT x Rf. "
     "Rf is the feedback resistor (user-selectable). "
     "The inverting input is a virtual ground. Output connects to ADC1 CH2 (PA1).",11,BODY,False,False,None),
    ("Impedance of DUT: |Z| = |V_ref| / |I_DUT| = (|V_ref| / |V_TIA|) x Rf.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("ADR4533BRZ — Precision Voltage Reference",13,NAVY,True,False,None),
    ("Provides 3.3 V with <25 uV/degC drift and <0.02% initial accuracy. "
     "Connected to STM32 VDDA and VREF+ pins to set the ADC full-scale voltage. "
     "Critical for measurement accuracy — noise on VREF adds directly to impedance error.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("MCP1700 — LDO Voltage Regulator",13,NAVY,True,False,None),
    ("USB-C VBUS (5 V) -> MCP1700 -> 3.3 V. Ultra-low quiescent current (1.6 uA typ.), "
     "250 mA output. Provides clean 3.3 V supply for the entire board. "
     "Input capacitor 1 uF, output capacitor 1 uF (per datasheet).",11,BODY,False,False,None),
],0.4,1.38,6.5,5.9)
# Right: signal chain block diagram
chain=[("AD9833 Output","~0.6 Vpp sine",BLUE),
       ("OPA2140 Buffer","Unity gain","Transimpedance amplifier is\nActual gain = 1",TEAL),
       ("Electrode (DUT)","V_exc / Z_dut = I",ORANGE),
       ("OPA2140 TIA","V = I x Rf",RED),
       ("ADC1 CH2 (PA1)","Digitize V_TIA",GREEN)]
cy10=1.42; bh10=0.75; gap10=0.2
for item in chain:
    if len(item)==3:
        name,sub,fill=item; lbox2(sl,name,sub,7.3,cy10,5.6,bh10,fill,tc=WHITE,tfs=11,sfs=9)
    else:
        name,sub,_,fill=item[0],item[1],item[2],item[2] if len(item)>3 else BLUE
    cy10+=bh10+gap10
    if cy10<6.8: cn(sl,10.1,cy10-gap10,10.1,cy10,color=BLUE,w=2)
# Fix: redo chain properly
cy10=1.42
chain2=[("AD9833 Output","~0.6 Vpp sine wave",BLUE),
        ("OPA2140 Buffer (Amp1A)","Unity gain, high-Z input -> Electrode+",TEAL),
        ("Electrode / DUT","Current I = V_exc / Z_dut flows through cell",ORANGE),
        ("OPA2140 TIA (Amp2A)","V_TIA = I x Rf  (transimpedance conversion)",RED),
        ("ADC1 CH2 — PA1","Digitizes V_TIA at 243 kHz, 256 samples",GREEN)]
# Clear previous and redraw
for i,(name,sub,fill) in enumerate(chain2):
    lbox2(sl,name,sub,7.3,cy10,5.6,bh10,fill,tc=WHITE,tfs=11,sfs=9)
    if i<len(chain2)-1:
        cn(sl,10.1,cy10+bh10,10.1,cy10+bh10+gap10,color=MUTED,w=2)
    cy10+=bh10+gap10
t(sl,"ADC1 CH1 (PA0) also captures V_exc reference simultaneously",7.3,cy10+0.05,5.6,0.45,fs=9.5,c=BLUE,italic=True)

# ═══ SLIDE 11 — Circuit Schematic (placeholder) ════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,NAVY)
t(sl,"Circuit Schematic",0.4,0.18,W-0.8,0.75,fs=32,c=WHITE,bold=True)
box(sl,0.4,0.92,W-0.8,0.04,TEAL)
box(sl,0.5,1.05,W-1.0,H-1.45,RGBColor(0x12,0x28,0x48),border=TEAL,bw=2)
t(sl,"[ Insert Circuit Schematic Here ]",0.5,H/2-0.5,W-1.0,0.75,fs=20,c=RGBColor(0x4C,0x8C,0xC4),align=PP_ALIGN.CENTER)
t(sl,"Right-click this box -> Change Picture -> From File -> select your schematic export (PNG/JPG)",0.5,H/2+0.35,W-1.0,0.5,fs=11,c=MUTED,italic=True,align=PP_ALIGN.CENTER)

# ═══ SLIDE 12 (NEW) — Excitation Signal Path ══════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Excitation Signal Path","From STM32 firmware command to electrode voltage — step by step")
txt(sl,[
    ("Step 1 — Firmware programs AD9833",13,NAVY,True,False,None),
    ("main.c calls AD9833_SetFrequency(freq). The function computes freqWord = freq x 268435456 / 25000000. "
     "FSYNC pulled LOW (PA4). Two 16-bit SPI writes sent: first word has D15=1, D14=0, carries bits[13:0] of freqWord. "
     "Second write carries bits[27:14]. FSYNC pulled HIGH to latch. AD9833 now outputs sine at requested frequency.",11,BODY,False,False,None),
    ("",5,BODY,False,False,None),
    ("Step 2 — Settle time",13,NAVY,True,False,None),
    ("HAL_Delay(SWEEP_DELAY_MS) gives the AD9833 output and the electrode/cell time to reach steady state "
     "before ADC acquisition begins. Default SWEEP_DELAY_MS = 100 ms.",11,BODY,False,False,None),
    ("",5,BODY,False,False,None),
    ("Step 3 — Signal buffered by OPA2140",13,NAVY,True,False,None),
    ("The AD9833 output (~0.6 Vpp centered near mid-supply) drives the non-inverting input of OPA2140 Amp1A "
     "wired as a unity-gain voltage follower. Output impedance drops to near zero, ensuring no amplitude "
     "error due to electrode loading.",11,BODY,False,False,None),
    ("",5,BODY,False,False,None),
    ("Step 4 — Signal applied to DUT",13,NAVY,True,False,None),
    ("The buffered sine wave appears at the Electrode+ terminal. Current flows through the "
     "electrochemical cell (DUT) with magnitude I = V_excitation / Z_DUT(f).",11,BODY,False,False,None),
],0.4,1.38,6.5,5.9)
# Right: schematic-style diagram
path_items=[("STM32 SPI1\n(PA4 FSYNC)","SPI command",NAVY),
            ("AD9833 DDS\n(25 MHz MCLK)","Sine wave output",BLUE),
            ("OPA2140 Amp1A\n(Unity-gain buffer)","Low-Z sine drive",TEAL),
            ("Electrode+\n(DUT terminal)","V_exc applied",ORANGE)]
cy12=1.5; bh12=1.0
for i,(name,label,fill) in enumerate(path_items):
    lbox2(sl,name,label,7.3,cy12,5.6,bh12,fill,tc=WHITE,tfs=12,sfs=10)
    if i<len(path_items)-1:
        cn(sl,10.1,cy12+bh12,10.1,cy12+bh12+0.22,color=BLUE,w=2.5)
    cy12+=bh12+0.22
t(sl,"Current I = V_exc / Z_DUT flows through cell",7.3,cy12+0.1,5.6,0.4,fs=10,c=ORANGE,italic=True,align=PP_ALIGN.CENTER)

# ═══ SLIDE 13 (NEW) — Measurement & Digitization Path ════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Measurement & Digitization Path","From electrode current to complex impedance — ADC + tone estimator")
txt(sl,[
    ("Step 1 — TIA converts current to voltage",13,NAVY,True,False,None),
    ("Cell current I flows into the inverting input of OPA2140 Amp2A (TIA). "
     "Feedback resistor Rf forces: V_TIA = -I x Rf = -(V_exc / Z_DUT) x Rf. "
     "V_TIA magnitude encodes |Z|; its phase relative to V_ref encodes the phase angle.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Step 2 — Dual-channel ADC acquisition via DMA",13,NAVY,True,False,None),
    ("ADC1 is configured in scan mode: CH1 (PA0 = V_ref) and CH2 (PA1 = V_TIA) sampled alternately. "
     "DMA transfers 512 half-words (256 per channel) to two separate buffers. "
     "Sample rate per channel = 72e6 / ((61.5+12.5) x 2) = 243 kHz.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Step 3 — Least-squares tone estimation",13,NAVY,True,False,None),
    ("For each channel, Estimate_Tone() projects the 256 samples onto cos(2*pi*f*t[i]) and sin(2*pi*f*t[i]) "
     "where t[i] = i / IMP_SAMPLE_RATE_HZ (actual hardware timing). "
     "Solving the 2x2 normal equations gives [a, b]. Magnitude = sqrt(a^2+b^2), Phase = atan2(b,a).",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Step 4 — Impedance calculation",13,NAVY,True,False,None),
    ("|Z| = (|V_ref| / |V_TIA|) x Rf  [Ohms]\n"
     "Phase = angle(V_TIA) - angle(V_ref)  [degrees]\n"
     "Re(Z) = |Z| x cos(phase)\n"
     "Im(Z) = |Z| x sin(phase)",11,BODY,False,False,None),
],0.4,1.38,6.5,5.9)
meas=[("Electrode- current","I = V_exc/Z_DUT",ORANGE),
      ("OPA2140 Amp2A TIA","V_TIA = -I x Rf",RED),
      ("ADC1 CH2 (PA1)","243 kHz, 256 samples",GREEN),
      ("Tone Estimator","a = proj(cos), b = proj(sin)",PURPLE),
      ("|Z|, Phase, Re, Im","BodePoint computed",NAVY)]
cy13=1.5; bh13=0.92
for i,(name,label,fill) in enumerate(meas):
    lbox2(sl,name,label,7.3,cy13,5.6,bh13,fill,tc=WHITE,tfs=12,sfs=10)
    if i<len(meas)-1: cn(sl,10.1,cy13+bh13,10.1,cy13+bh13+0.18,color=TEAL,w=2.5)
    cy13+=bh13+0.18

# ═══ SLIDE 14 — PCB Layout (placeholder) ══════════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,NAVY)
t(sl,"PCB Layout",0.4,0.18,W-0.8,0.75,fs=32,c=WHITE,bold=True)
box(sl,0.4,0.92,W-0.8,0.04,TEAL)
box(sl,0.5,1.05,W-1.0,H-1.45,RGBColor(0x12,0x28,0x48),border=TEAL,bw=2)
t(sl,"[ Insert PCB Layout Image Here ]",0.5,H/2-0.5,W-1.0,0.75,fs=20,c=RGBColor(0x4C,0x8C,0xC4),align=PP_ALIGN.CENTER)
t(sl,"Right-click -> Change Picture -> From File -> select the PCB screenshot from KiCad/Altium",0.5,H/2+0.35,W-1.0,0.5,fs=11,c=MUTED,italic=True,align=PP_ALIGN.CENTER)

# ═══ SLIDE 15 (NEW) — PCB Design Considerations ══════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"PCB Design Considerations","Key design decisions that affect measurement accuracy and signal integrity")
considerations=[
    ("Analog/Digital Ground Plane",
     "Analog and digital sections share a common ground but with careful layout to avoid digital switching noise "
     "coupling into the sensitive TIA output trace. Decoupling capacitors placed directly at each IC power pin.",NAVY),
    ("Short TIA Output Trace",
     "The trace from OPA2140 Amp2A output to ADC1 CH2 (PA1) is kept as short as possible to minimize parasitic "
     "capacitance, which would add phase error and bandwidth limitation at high frequencies.",BLUE),
    ("VREF+ Filtering",
     "ADR4533BRZ output is decoupled with 100 nF + 1 uF capacitors directly at the STM32 VREF+ pin. "
     "Any ripple on the voltage reference appears directly as gain error in all ADC readings.",TEAL),
    ("SPI Bus Routing",
     "SPI1 clock, MOSI, and MISO traces are routed with matched length where practical. Both AD9833 (CS=PA4) "
     "and W25Q32 (CS=PB12) share the three data lines — firmware ensures mutual exclusion.",GREEN),
    ("Crystal Placement",
     "The 8 MHz STM32 crystal (ABM3B) is placed very close to the HSE pins with guard traces. "
     "The 25 MHz AD9833 crystal (ASE-25.000MHz) is placed directly adjacent to the AD9833 MCLK pin.",ORANGE),
    ("USB-C Power Path",
     "USB-C VBUS enters through a 5.1 kOhm CC resistor for DFP detection, then into the MCP1700 LDO. "
     "Input bulk capacitor (10 uF) absorbs USB cable inductance transients before the LDO.",PURPLE),
]
rw15=(W-0.9)/2; rh15=1.4
for i,(title,desc,fill) in enumerate(considerations):
    col=i%2; row=i//2
    sx=0.4+col*(rw15+0.1); sy=1.4+row*(rh15+0.1)
    lbox2(sl,title,desc,sx,sy,rw15,rh15,fill,tc=WHITE,tfs=11,sfs=9.5)

# ═══ SLIDE 16 (NEW) — Component Placement & Routing Strategy ═════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Component Placement & Routing Strategy","Where each IC sits on the PCB and why")
txt(sl,[
    ("STM32F303CCT6 — Center of Board",13,NAVY,True,False,None),
    ("Placed centrally so all peripheral connections (SPI, UART, ADC, GPIO) radiate outward with minimal trace length.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("AD9833 DDS — Near STM32 SPI Pins",13,BLUE,True,False,None),
    ("Short SPI traces reduce impedance mismatch at 36 MHz SPI clock. "
     "FSYNC trace (PA4) goes directly to AD9833 SYNC pin.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("OPA2140 Pair — Near Electrode Connector",13,TEAL,True,False,None),
    ("Minimises trace length for the high-impedance TIA input (inverting node) and excitation output, "
     "reducing parasitic capacitance pickup from nearby digital traces.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("ADR4533BRZ — Adjacent to VREF+ Pin",13,ORANGE,True,False,None),
    ("Direct routing from ADR4533 output to VREF+ with 100nF decoupling cap on the same net, "
     "with 1uF bulk cap to ground.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("W25Q32 Flash — On SPI1 Bus, Opposite Side from AD9833",13,GREEN,True,False,None),
    ("Separate chip-select (PB12) allows the firmware to select flash vs DDS independently. "
     "CS line pulled HIGH at power-on via pull-up resistor (PB12 pull-up).",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("CP2102 & USB-C — Board Edge",13,PURPLE,True,False,None),
    ("USB connector on board edge with ESD protection diodes. CP2102 placed adjacent to reduce "
     "UART trace length to STM32 PA9/PA10.",11,BODY,False,False,None),
],0.4,1.38,6.5,5.9)
# Right: placement diagram (abstract top-view of PCB)
box(sl,7.0,1.38,6.0,5.9,RGBColor(0x0A,0x2A,0x10),border=TEAL,bw=2)
t(sl,"PCB Top View (Abstract)",7.1,1.45,5.8,0.35,fs=11,c=TEAL,bold=True,align=PP_ALIGN.CENTER)
placements=[("STM32F303",9.5,3.5,1.6,1.0,NAVY),("AD9833",7.5,2.0,1.4,0.7,BLUE),
            ("OPA2140x2",7.5,4.5,1.4,0.7,TEAL),("ADR4533",11.0,2.0,1.3,0.7,ORANGE),
            ("W25Q32",11.0,4.5,1.3,0.7,GREEN),("CP2102",11.0,3.4,1.3,0.7,RED),
            ("MCP1700",7.5,3.3,1.2,0.6,PURPLE),("8MHz Xtal",9.0,2.0,1.0,0.55,BROWN),
            ("25MHz Xtal",8.3,4.1,1.1,0.55,AMBER)]
for name,px,py,pw,ph,fill in placements:
    lbox(sl,name,px,py,pw,ph,fill,tc=WHITE,fs=8,bold=True)
lbox(sl,"Electrode\nConnector",7.2,5.5,1.2,0.6,RGBColor(0x37,0x47,0x4F),tc=WHITE,fs=8)
lbox(sl,"USB-C",11.6,5.5,1.1,0.6,RGBColor(0x37,0x47,0x4F),tc=WHITE,fs=8)

# ═══ SLIDE 17 — Firmware & Software ══════════════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Firmware & Software","STM32 HAL-based C firmware  |  Python 3 desktop application")
txt(sl,[
    ("STM32 Firmware (C / HAL)",13,NAVY,True,False,None),
    ("The firmware initialises SPI1, configures dual-channel ADC1 with DMA (486 kHz effective rate, "
     "256 samples), and enters a UART command loop at 115 200 baud. On receiving START it iterates "
     "through the user-specified frequency sweep calling AD9833 to set each frequency.",10,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("At each frequency step, 512 ADC samples (256 x ref + 256 x TIA) are collected via DMA. "
     "A least-squares 2x2 sinusoidal tone estimator projects onto the exact measurement frequency "
     "yielding true complex voltage for both channels. Z = (V_ref/V_sig) x Rf.",10,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Each BodePoint (|Z|, phase, Re, Im) is transmitted over UART and simultaneously written to "
     "W25Q32 SPI flash. DUMP_FLASH replays stored records to the GUI.",10,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Python GUI (PyQt5)",13,NAVY,True,False,None),
    ("A background QThread reads UART lines at 115200 baud. The main thread parses "
     "DATA/STATUS/FLASH responses and updates three live Matplotlib canvases (Bode |Z|, phase, Nyquist). "
     "ConfigPanel manages sweep parameters and displays live measurement state.",10,BODY,False,False,None),
],0.4,1.38,6.5,5.9)
module_groups=[("Firmware Modules",NAVY,[
    ("main.c","Sweep orchestration, flash init, status reporting"),
    ("impedance.c","Least-squares tone estimator, BodePoint calc"),
    ("ad9833.c","DDS SPI driver: freq word -> 2 x SPI writes"),
    ("w25q32.c","Low-level SPI flash: page write, sector erase"),
    ("flash_log.c","Structured log: BeginSweep, WritePoint, Dump"),
    ("uart_comm.c","ASCII command parser: START/STOP/STATUS/FLASH"),
]),("Python Modules",TEAL,[
    ("main.py","QMainWindow; toolbar; UART line parser"),
    ("serial_worker.py","QThread; COM port open/close; line reader"),
    ("data_model.py","ImpedancePoint list; CSV export"),
    ("plot_widget.py","3x Matplotlib FigureCanvas"),
    ("config_panel.py","Sweep settings; live status; last point"),
])]
rx,ry=7.2,1.38; rw=W-rx-0.4
for grp_title,grp_color,items in module_groups:
    lbox(sl,grp_title,rx,ry,rw,0.44,grp_color,tc=WHITE,fs=12,bold=True); ry+=0.44
    for name,desc in items:
        lbox2(sl,name,desc,rx,ry,rw,0.54,CARD,tc=BODY,tfs=10,sfs=8.5,border=grp_color); ry+=0.55
    ry+=0.12

# ═══ SLIDE 18 (NEW) — ADC Configuration & DMA Transfer ══════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"ADC Configuration & DMA Transfer","How the STM32 ADC1 is set up for simultaneous dual-channel acquisition")
txt(sl,[
    ("ADC1 Configuration",13,NAVY,True,False,None),
    ("Mode: Scan conversion — ADC1 cycles through CH1 (PA0) and CH2 (PA1) in sequence.",11,BODY,False,False,None),
    ("Resolution: 12-bit (0–4095), right-aligned.",11,BODY,False,False,None),
    ("Sample time: 61.5 ADC clock cycles per channel.",11,BODY,False,False,None),
    ("Conversion time: 12.5 ADC clock cycles per channel.",11,BODY,False,False,None),
    ("ADC clock = 72 MHz / 1 = 72 MHz (AHB, no prescaler in this config).",11,BODY,False,False,None),
    ("Time per sample = (61.5 + 12.5) / 72e6 = 1.028 us.",11,BODY,False,False,None),
    ("Effective rate per channel = 1 / (2 x 1.028 us) = 486 kHz total / 2 = 243 kHz.",11,BODY,False,False,None),
    ("",5,BODY,False,False,None),
    ("DMA Configuration",13,NAVY,True,False,None),
    ("DMA1 Channel 1 is mapped to ADC1. Transfer size = 512 half-words (uint16_t).",11,BODY,False,False,None),
    ("Memory layout: [CH1[0], CH2[0], CH1[1], CH2[1], ..., CH1[255], CH2[255]]",11,BODY,False,False,None),
    ("After DMA complete interrupt, firmware separates interleaved samples into two float arrays: "
     "ref[256] and sig[256] before calling Estimate_Tone().",11,BODY,False,False,None),
    ("",5,BODY,False,False,None),
    ("Timing per frequency step",13,NAVY,True,False,None),
    ("Settle delay: SWEEP_DELAY_MS (default 100 ms). DMA acquisition: 512 / 486k = 1.05 ms. "
     "Tone estimation: < 1 ms. Flash write: ~2 ms. Total: ~104 ms/step.",11,BODY,False,False,None),
],0.4,1.38,6.5,5.9)
# Right: DMA buffer layout diagram
rbox(sl,7.2,1.5,5.8,5.5,WHITE,border=BLUE,bw=1)
t(sl,"DMA Buffer Layout (512 half-words)",7.4,1.6,5.4,0.35,fs=12,c=NAVY,bold=True,align=PP_ALIGN.CENTER)
cells=[("CH1[0]\nRef",BLUE),("CH2[0]\nSig",TEAL),("CH1[1]\nRef",BLUE),("CH2[1]\nSig",TEAL),
       ("...",MUTED),("...",MUTED),("CH1[255]\nRef",BLUE),("CH2[255]\nSig",TEAL)]
cx18=7.3; cy18=2.1; cw18=1.35; ch18=0.75
for i,(label,fill) in enumerate(cells):
    col=i%4; row=i//4
    lbox(sl,label,cx18+col*cw18,cy18+row*(ch18+0.08),cw18-0.04,ch18,fill,tc=WHITE,fs=8.5,bold=True)
t(sl,"After DMA complete -> separate into ref[] and sig[]",7.3,3.85,5.7,0.4,fs=10,c=NAVY,bold=True,align=PP_ALIGN.CENTER)
lbox(sl,"ref[0..255]  =  even indices",7.3,4.32,2.75,0.65,BLUE,tc=WHITE,fs=10)
lbox(sl,"sig[0..255]  =  odd indices",10.2,4.32,2.75,0.65,TEAL,tc=WHITE,fs=10)
cn(sl,10.1,5.05,10.1,5.3,color=MUTED,w=1.5)
lbox(sl,"Estimate_Tone(ref, freq) -> V_ref (mag, phase)",7.3,5.32,2.75,0.65,BLUE,tc=WHITE,fs=9)
lbox(sl,"Estimate_Tone(sig, freq) -> V_sig (mag, phase)",10.2,5.32,2.75,0.65,TEAL,tc=WHITE,fs=9)

# ═══ SLIDE 19 (NEW) — Least-Squares Tone Estimator ══════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Least-Squares Tone Estimator Algorithm","Why naive DFT fails and how the LS projector solves it")
txt(sl,[
    ("Problem with Naive DFT",13,RED,True,False,None),
    ("Standard DFT assumes the signal frequency is exactly k/N (a normalized bin). "
     "For EIS, the measurement frequency f is arbitrary and typically not a bin. "
     "This causes spectral leakage: the computed magnitude is too low and the phase is wrong.",11,BODY,False,False,None),
    ("",5,BODY,False,False,None),
    ("Least-Squares Projection Solution",13,NAVY,True,False,None),
    ("Model: x[i] = a*cos(2*pi*f*t[i]) + b*sin(2*pi*f*t[i])  where t[i] = i / IMP_SAMPLE_RATE_HZ",11,BODY,False,False,None),
    ("Build sums:  Scc = sum(cos^2),  Sss = sum(sin^2),  Scs = sum(cos*sin)",11,BODY,False,False,None),
    ("            Scx = sum(x[i]*cos(2pi*f*t[i])),  Ssx = sum(x[i]*sin(2pi*f*t[i]))",11,BODY,False,False,None),
    ("Solve 2x2 system:  [Scc Scs] [a]   [Scx]",11,BODY,False,False,None),
    ("                   [Scs Sss] [b] = [Ssx]",11,BODY,False,False,None),
    ("Determinant D = Scc*Sss - Scs^2",11,BODY,False,False,None),
    ("Solution:  a = (Sss*Scx - Scs*Ssx) / D",11,BODY,False,False,None),
    ("           b = (Scc*Ssx - Scs*Scx) / D",11,BODY,False,False,None),
    ("",5,BODY,False,False,None),
    ("Result",13,NAVY,True,False,None),
    ("magnitude = sqrt(a^2 + b^2)     [proportional to channel voltage amplitude]",11,BODY,False,False,None),
    ("phaseRad  = atan2(b, a)          [phase of the signal relative to t=0]",11,BODY,False,False,None),
    ("Z = (ref.magnitude / sig.magnitude) * Rf   [Ohms]",11,BODY,False,False,None),
    ("phi = (sig.phaseRad - ref.phaseRad) * 180/pi  [degrees]",11,BODY,False,False,None),
],0.4,1.38,6.5,5.9)
# Right: comparison diagram
rbox(sl,7.2,1.45,5.8,2.6,WHITE,border=RED,bw=1)
t(sl,"Naive DFT (wrong at non-bin freq)",7.4,1.52,5.4,0.35,fs=11,c=RED,bold=True,align=PP_ALIGN.CENTER)
t(sl,"Projects onto cos(2pi*k/N*i) — correct only if f = k*Fs/N\nAt other frequencies: amplitude error + phase error\nCan give |Z| error >10% at typical EIS frequencies",7.4,1.9,5.4,1.0,fs=10,c=BODY)
rbox(sl,7.2,4.2,5.8,2.6,WHITE,border=GREEN,bw=1)
t(sl,"LS Tone Estimator (correct at any freq)",7.4,4.27,5.4,0.35,fs=11,c=GREEN,bold=True,align=PP_ALIGN.CENTER)
t(sl,"Projects onto cos(2pi*f*i/Fs) using real hardware timing\nNo spectral leakage assumption\nAmplitude and phase accurate at any f within Nyquist\nCost: ~4*N multiplications per channel",7.4,4.65,5.4,1.2,fs=10,c=BODY)
box(sl,9.6,4.12,0.8,0.12,GREEN)
t(sl,"CORRECT",9.6,3.85,1.4,0.28,fs=9,c=GREEN,bold=True,align=PP_ALIGN.CENTER)

# ═══ SLIDE 20 — Software Data Flow ═══════════════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Software: Data Flow & Module Interactions","Left: command/data pipeline  |  Right: module dependency graph")
fc_x,fc_w,fc_h,fc_gap=0.35,3.8,0.54,0.16
steps20=[("GUI: user clicks Start Sweep",BLUE),("SerialWorker: send 'START' over UART",TEAL),
          ("STM32 parser: enter sweep loop",NAVY),("AD9833: program frequency (SPI x 2)",ORANGE),
          ("ADC DMA: capture 256 x 2 samples",GREEN),("Tone estimator: solve 2x2 LS at f",PURPLE),
          ("Compute |Z|, Phase, Re(Z), Im(Z)",NAVY),("UART TX: DATA,f,|Z|,phase + Flash write",RED),
          ("GUI parser: update Bode + Nyquist",BLUE)]
cy=1.38
for i,(label,fill) in enumerate(steps20):
    lbox(sl,label,fc_x,cy,fc_w,fc_h,fill,tc=WHITE,fs=9.5,bold=False)
    if i<len(steps20)-1: cn(sl,fc_x+fc_w/2,cy+fc_h,fc_x+fc_w/2,cy+fc_h+fc_gap+0.01,color=BLUE,w=1.8)
    cy+=fc_h+fc_gap
mw2=2.2; mh2=0.58
row_y=[1.38,2.55,3.72,4.89]
mods20={"Python GUI":(4.65,row_y[0],BLUE),"uart_comm.c":(7.35,row_y[0],TEAL),
        "main.c":(5.55,row_y[1],NAVY),"flash_log.c":(8.15,row_y[1],GREEN),
        "ad9833.c":(4.30,row_y[2],ORANGE),"impedance.c":(6.80,row_y[2],PURPLE),
        "w25q32.c":(9.30,row_y[2],RED),"AD9833 (HW)":(4.30,row_y[3],AMBER),
        "ADC1/DMA":(6.80,row_y[3],BROWN),"W25Q32 (HW)":(9.30,row_y[3],SLATE)}
for name,(mx2,my2,fill) in mods20.items():
    lbox(sl,name,mx2,my2,mw2,mh2,fill,tc=WHITE,fs=9,bold=True)
def mc20(n): mx2,my2,_=mods20[n]; return mx2+mw2/2,my2+mh2/2
for a,b in [("Python GUI","uart_comm.c"),("uart_comm.c","main.c"),("main.c","flash_log.c"),
            ("main.c","ad9833.c"),("main.c","impedance.c"),("flash_log.c","w25q32.c"),
            ("impedance.c","ADC1/DMA"),("ad9833.c","AD9833 (HW)"),("w25q32.c","W25Q32 (HW)")]:
    ax,ay=mc20(a); bx,by=mc20(b); cn(sl,ax,ay,bx,by,color=MUTED,w=1.0)

# ═══ SLIDE 21 (NEW) — Python GUI Data Pipeline ════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Python GUI — Data Pipeline Detail","How incoming UART bytes become live plot updates")
txt(sl,[
    ("1. SerialWorker.run() — Background Thread",13,NAVY,True,False,None),
    ("Reads available bytes from PySerial port (non-blocking, 10 ms sleep if empty). "
     "Appends bytes to a buffer. On each \\n found, strips the line and emits line_received(str) signal. "
     "Thread-safe: Qt signals cross threads automatically.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("2. MainWindow._on_line_received() — Main Thread",13,NAVY,True,False,None),
    ("Routes by prefix: DATA -> _parse_data(), FLASH_DATA -> _parse_flash_data(), "
     "STATUS -> _parse_status(), FLASH_STATUS -> _parse_flash_status(), "
     "SWEEP,BEGIN / SWEEP,DONE -> state update, ERROR -> red console log.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("3. _parse_data() — For each DATA line",13,NAVY,True,False,None),
    ("Splits CSV: freq, |Z|, phase. Computes Re = |Z|*cos(phase), Im = |Z|*sin(phase). "
     "Calls DataModel.add_point(). Updates points counter. "
     "Calls ConfigPanel.update_last_point() and update_status(). "
     "Calls PlotWidget.update_plots() with full arrays.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("4. PlotWidget.update_plots() — Redraw",13,NAVY,True,False,None),
    ("Clears all three axes, replots with current full dataset, calls canvas.draw(). "
     "Bode |Z| is log-log, Bode phase is log-x linear, Nyquist plots Re vs -Im with equal aspect.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("5. Export",13,NAVY,True,False,None),
    ("DataModel.export_csv(path) writes a header row then one CSV row per ImpedancePoint. "
     "Works on live or flash-dumped data identically.",11,BODY,False,False,None),
],0.4,1.38,6.5,5.9)
# Right: pipeline flowchart
pipe21=[("Serial port bytes","PySerial read",SLATE),
        ("SerialWorker buffer","Byte accumulation",SLATE),
        ("line_received(str)","Qt signal across threads",BLUE),
        ("_on_line_received()","Route by prefix",NAVY),
        ("_parse_data()","Split CSV, compute Re/Im",TEAL),
        ("DataModel","add_point()",GREEN),
        ("PlotWidget","update_plots()",ORANGE),
        ("ConfigPanel","update_last_point()",PURPLE)]
cy21=1.45; bh21=0.65
for i,(name,sub,fill) in enumerate(pipe21):
    lbox2(sl,name,sub,7.2,cy21,5.8,bh21,fill,tc=WHITE,tfs=11,sfs=9)
    if i<len(pipe21)-1: cn(sl,10.1,cy21+bh21,10.1,cy21+bh21+0.14,color=BLUE,w=2)
    cy21+=bh21+0.14

# ═══ SLIDE 22 (NEW) — Firmware UART Command Table ════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Firmware UART Command Reference","Complete ASCII protocol — commands sent by GUI and responses from STM32")
commands=[
    ("START","Begin frequency sweep. Calls FlashLog_BeginSweep() then iterates sweep loop.","SWEEP,BEGIN -> DATA,f,|Z|,phase (per step) -> SWEEP,DONE",BLUE),
    ("STOP","Sets measurement_running=0. Current step completes before stopping.","(no explicit response; next STATUS will show IDLE)",TEAL),
    ("STATUS","Report current operating state and sweep parameters.","STATUS,IDLE|RUNNING,START=n,STOP=n,STEP=n,RF=n,FLASH_COUNT=n",NAVY),
    ("SET_START_FREQ,n","Set sweep start frequency to n Hz.","+OK or no response",GREEN),
    ("SET_STOP_FREQ,n","Set sweep stop frequency to n Hz.","+OK or no response",GREEN),
    ("SET_STEP_FREQ,n","Set sweep frequency step to n Hz.","+OK or no response",GREEN),
    ("SET_RF,n.n","Set TIA feedback resistor value to n.n Ohms.","+OK or no response",GREEN),
    ("FLASH_STATUS","Query flash log record count and capacity.","FLASH_STATUS,COUNT=n,CAPACITY=m",ORANGE),
    ("DUMP_FLASH","Replay all stored flash records to UART.","FLASH_DATA,f,|Z|x100,phasex100,Rex100,Imx100 (per record)",PURPLE),
    ("ERASE_FLASH","Erase the flash log region (4 sectors x 4 KB = 16 KB).","FLASH_ERASE_OK or FLASH_ERASE_ERROR",RED),
]
rh22=0.52; rw_cmd=2.5; rw_desc=4.5; rw_resp=5.5
lbox(sl,"Command",0.4,1.38,rw_cmd,0.45,NAVY,tc=WHITE,fs=11,bold=True)
lbox(sl,"Action",0.4+rw_cmd,1.38,rw_desc,0.45,NAVY,tc=WHITE,fs=11,bold=True)
lbox(sl,"Response",0.4+rw_cmd+rw_desc,1.38,rw_resp,0.45,NAVY,tc=WHITE,fs=11,bold=True)
cy22=1.38+0.45
for cmd,action,response,fill in commands:
    lbox(sl,cmd,0.4,cy22,rw_cmd,rh22,fill,tc=WHITE,fs=9,bold=True)
    lbox(sl,action,0.4+rw_cmd,cy22,rw_desc,rh22,CARD,tc=BODY,fs=8.5,bold=False,align=PP_ALIGN.LEFT)
    lbox(sl,response,0.4+rw_cmd+rw_desc,cy22,rw_resp,rh22,LCARD,tc=BODY,fs=8,bold=False,align=PP_ALIGN.LEFT)
    cy22+=rh22

# ═══ SLIDE 23 — Hardware Interaction Diagram ══════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Hardware Component Interactions","All peripherals connected to STM32F303CCT6")
stm_x,stm_y,stm_w,stm_h=5.0,3.0,3.4,1.15
lbox(sl,"STM32F303CCT6\n72 MHz Cortex-M4",stm_x,stm_y,stm_w,stm_h,NAVY,tc=WHITE,fs=13,bold=True)
stm_cx=stm_x+stm_w/2; stm_cy=stm_y+stm_h/2
hw_pw,hw_ph=2.1,0.72
perif23=[(0.2,1.1,"AD9833\nDDS",BLUE,"SPI1  PA4 FSYNC"),
          (4.2,0.5,"8 MHz Crystal\n(HSE)",ORANGE,"PF0/PF1"),
          (7.5,1.1,"W25Q32 Flash\n32 Mbit",GREEN,"SPI1  PB12 CS"),
          (10.8,1.1,"CP2102\nUSB-UART",TEAL,"USART1 PA9/PA10"),
          (0.2,3.5,"MCP1700\nLDO 3.3V",PURPLE,"3.3V power rail"),
          (0.2,5.2,"OPA2140\nExcit. Buffer",AMBER,"DAC Out->Electrode"),
          (10.8,3.5,"ADR4533\n3.3V Ref",RED,"VDDA/VREF+"),
          (10.8,5.2,"OPA2140\nTIA (I->V)",ORANGE,"ADC1 CH2 PA1"),
          (4.2,5.6,"ADC1 (DMA)\n2ch 256 smp",BROWN,"PA0  PA1  486kHz"),
          (7.5,5.6,"25 MHz Crystal\n(AD9833 MCLK)",CYAN,"CLKIN")]
for px2,py2,label,fill,bus_lbl in perif23:
    lbox(sl,label,px2,py2,hw_pw,hw_ph,fill,tc=WHITE,fs=9,bold=True)
    pcx=px2+hw_pw/2; pcy=py2+hw_ph/2
    cn(sl,pcx,pcy,stm_cx,stm_cy,color=MUTED,w=1.0)
    t(sl,bus_lbl,(pcx+stm_cx)/2-0.65,(pcy+stm_cy)/2-0.16,1.5,0.3,fs=7,c=MUTED,italic=True)

# ═══ SLIDE 24 (NEW) — SPI Bus Sharing ════════════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"SPI Bus Sharing: AD9833 + W25Q32","How two SPI devices share one bus with independent chip-selects")
txt(sl,[
    ("SPI1 Bus Configuration",13,NAVY,True,False,None),
    ("SPI1 pins: PA5=SCK, PA6=MISO, PA7=MOSI. Hardware SPI at 36 MHz (72 MHz / prescaler 2). "
     "Mode 2 (CPOL=1, CPHA=0) compatible with both AD9833 and W25Q32.",11,BODY,False,False,None),
    ("",5,BODY,False,False,None),
    ("AD9833 Chip-Select: PA4 (FSYNC)",13,BLUE,True,False,None),
    ("Controlled as GPIO output. FSYNC_LOW() macro pulls PA4 low before SPI transfer. "
     "FSYNC_HIGH() pulls PA4 high after transfer to latch the frequency word. "
     "AD9833 is write-only — no MISO data is returned.",11,BODY,False,False,None),
    ("",5,BODY,False,False,None),
    ("W25Q32 Chip-Select: PB12",13,GREEN,True,False,None),
    ("Controlled as GPIO output. W25Q32_CS_LOW() pulls PB12 low, W25Q32_CS_HIGH() releases. "
     "W25Q32 supports full-duplex: MOSI carries commands/data, MISO returns read data. "
     "Software pull-up on PB12 at boot to keep flash deselected during STM32 startup.",11,BODY,False,False,None),
    ("",5,BODY,False,False,None),
    ("Mutual Exclusion",13,NAVY,True,False,None),
    ("Firmware never asserts both CS lines simultaneously. AD9833 operations (AD9833_SetFrequency) "
     "complete before any W25Q32 operation (FlashLog_WritePoint) begins. No RTOS mutex needed "
     "since the single sweep thread owns both devices.",11,BODY,False,False,None),
    ("",5,BODY,False,False,None),
    ("Bus Sharing Protocol",13,TEAL,True,False,None),
    ("1. Assert AD9833 FSYNC (PA4 LOW)  2. SPI transfer  3. Deassert FSYNC (PA4 HIGH)\n"
     "4. Assert W25Q32 CS (PB12 LOW)    5. SPI transfer  6. Deassert CS (PB12 HIGH)",11,BODY,False,False,None),
],0.4,1.38,6.5,5.9)
# Right: bus diagram
lbox(sl,"SPI1 Bus\nPA5=SCK  PA6=MISO  PA7=MOSI",7.2,1.5,5.8,0.85,NAVY,tc=WHITE,fs=11,bold=True)
cn(sl,9.4,2.35,9.4,2.75,color=NAVY,w=2)
lbox(sl,"AD9833 DDS\n(PA4=FSYNC, write-only)",7.2,2.75,2.7,0.85,BLUE,tc=WHITE,fs=10)
lbox(sl,"W25Q32 Flash\n(PB12=CS, read/write)",10.1,2.75,2.9,0.85,GREEN,tc=WHITE,fs=10)
t(sl,"PA4 LOW -> transfer -> PA4 HIGH",7.2,3.7,2.7,0.45,fs=9,c=BLUE,italic=True,align=PP_ALIGN.CENTER)
t(sl,"PB12 LOW -> transfer -> PB12 HIGH",10.1,3.7,2.9,0.45,fs=9,c=GREEN,italic=True,align=PP_ALIGN.CENTER)
lbox(sl,"Never both CS lines active simultaneously",7.2,4.25,5.8,0.55,RED,tc=WHITE,fs=11,bold=True)
# Show timing
rbox(sl,7.2,5.0,5.8,2.1,LCARD,border=NAVY,bw=1)
t(sl,"Timing sequence (per frequency step):",7.4,5.08,5.4,0.32,fs=10,c=NAVY,bold=True)
seq=[("AD9833_SetFrequency(f)","PA4 low -> 2x SPI word -> PA4 high",BLUE),
     ("HAL_Delay(100 ms)","Signal settle time",MUTED),
     ("ADC DMA acquisition","No SPI activity",GREEN),
     ("FlashLog_WritePoint()","PB12 low -> SPI R/W -> PB12 high",GREEN)]
sy=5.45
for name,desc,c in seq:
    t(sl,f"{name}: {desc}",7.35,sy,5.5,0.35,fs=8.5,c=c)
    sy+=0.37

# ═══ SLIDE 25 (NEW) — Analog Signal Chain & Signal Levels ════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Analog Signal Chain & Signal Levels","Voltage and current at each stage from DDS to ADC")
# Left: detailed text
txt(sl,[
    ("Stage 1: AD9833 Output",13,BLUE,True,False,None),
    ("Sine wave, ~0.6 Vpp, centered around ~1.65 V (mid-supply). "
     "Frequency set by SPI command. Output impedance ~200 Ohms.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Stage 2: OPA2140 Excitation Buffer",13,TEAL,True,False,None),
    ("Unity-gain voltage follower. Input impedance > 1 TOhm (FET input). "
     "Output impedance < 1 Ohm. Drives Electrode+ terminal with full fidelity.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Stage 3: Electrode / DUT",13,ORANGE,True,False,None),
    ("Current I(f) = V_exc / Z_DUT(f) flows from Electrode+ through the electrochemical cell. "
     "Z_DUT varies with frequency — this is what we are measuring.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Stage 4: OPA2140 TIA",13,RED,True,False,None),
    ("Inverting TIA configuration. Virtual ground at inverting input (Electrode-). "
     "V_TIA = -I x Rf = -(V_exc / Z_DUT) x Rf. "
     "Output range must stay within ADC input range (0 to 3.3V).",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Stage 5: ADC Reference Channel",13,NAVY,True,False,None),
    ("PA0 (ADC1 CH1) samples V_exc directly. This is the reference for phase and amplitude. "
     "Same ADC, same timing -> no systematic phase offset between channels.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Choosing Rf",13,GREEN,True,False,None),
    ("Select Rf so |V_TIA| at the expected |Z_DUT| fills ~50-70% of ADC range for best SNR. "
     "Rf = 10 kOhm is a good starting point for 1-100 kOhm cells.",11,BODY,False,False,None),
],0.4,1.38,6.5,5.9)
# Right: signal level bar chart (drawn with shapes)
rbox(sl,7.0,1.45,6.0,5.8,WHITE,border=BLUE,bw=1)
t(sl,"Signal Levels at Each Stage",7.2,1.55,5.6,0.35,fs=11,c=NAVY,bold=True,align=PP_ALIGN.CENTER)
stages25=[("AD9833 Out","~0.6 Vpp",BLUE,3.0),("Buffer Out","~0.6 Vpp",TEAL,3.0),
          ("V_exc at DUT","~0.6 Vpp",ORANGE,3.0),("I thru DUT","depends on Z",RED,1.5),
          ("V_TIA","I x Rf",RED,2.0),("ADC CH2","0 to 3.3V range",GREEN,3.3),
          ("ADC CH1 (Ref)","~0.6 Vpp",NAVY,3.0)]
bar_x=7.2; bar_max_w=5.6; bar_h=0.52; bar_y=2.05; max_val=3.3
for name,val,fill,bar_len in stages25:
    bw25=bar_max_w*bar_len/max_val
    box(sl,bar_x,bar_y,bw25,bar_h,fill)
    t(sl,f"{name}: {val}",bar_x+0.06,bar_y+0.12,bw25+1.5,0.3,fs=8.5,c=BODY,bold=False)
    bar_y+=bar_h+0.1

# ═══ SLIDE 26 — System Integration ═══════════════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"System Integration","How Python GUI, STM32 firmware, and hardware signal chain cooperate")
col_w3=3.8; col_h_hdr=0.72; col_gap=0.55; cy_start=1.42
col_xs=[0.4,0.4+col_w3+col_gap,0.4+2*(col_w3+col_gap)]
lbox(sl,"Python GUI\n(Laptop/PC)",col_xs[0],cy_start,col_w3,col_h_hdr,BLUE,tc=WHITE,fs=14,bold=True)
lbox(sl,"STM32 Firmware\n(Embedded C)",col_xs[1],cy_start,col_w3,col_h_hdr,NAVY,tc=WHITE,fs=14,bold=True)
lbox(sl,"Hardware\n(Signal Chain)",col_xs[2],cy_start,col_w3,col_h_hdr,TEAL,tc=WHITE,fs=14,bold=True)
gui_items=["SerialWorker (QThread)","UART line parser","DataModel (store)","PlotWidget (Bode+Nyquist)","ConfigPanel (settings)","CSV export / flash dump"]
fw_items=["uart_comm.c (parser)","main.c (sweep loop)","impedance.c (tone est.)","ad9833.c (DDS driver)","flash_log.c (log)","w25q32.c (flash R/W)"]
hw_items=["AD9833 DDS (sine gen.)","OPA2140 excit. buffer","Electrode / DUT","OPA2140 TIA","ADC1 with DMA (2ch)","W25Q32 SPI flash"]
item_h=0.46
for i,(gi,fi,hi) in enumerate(zip(gui_items,fw_items,hw_items)):
    iy=cy_start+col_h_hdr+i*(item_h+0.04)
    lbox(sl,gi,col_xs[0],iy,col_w3,item_h,LCARD,tc=BODY,fs=9.5,bold=False,border=BLUE,bw=0.5)
    lbox(sl,fi,col_xs[1],iy,col_w3,item_h,LCARD,tc=BODY,fs=9.5,bold=False,border=NAVY,bw=0.5)
    lbox(sl,hi,col_xs[2],iy,col_w3,item_h,LCARD,tc=BODY,fs=9.5,bold=False,border=TEAL,bw=0.5)
arrow_y=cy_start+col_h_hdr/2; ret_y=arrow_y+0.44
cn(sl,col_xs[0]+col_w3,arrow_y,col_xs[1],arrow_y,color=BLUE,w=2.2)
cn(sl,col_xs[1],ret_y,col_xs[0]+col_w3,ret_y,color=NAVY,w=2.2)
cn(sl,col_xs[1]+col_w3,arrow_y,col_xs[2],arrow_y,color=TEAL,w=2.2)
t(sl,"UART cmds",col_xs[0]+col_w3+0.04,arrow_y-0.42,col_gap-0.06,0.38,fs=8,c=BLUE,italic=True,align=PP_ALIGN.CENTER)
t(sl,"DATA/STATUS",col_xs[0]+col_w3+0.04,ret_y+0.04,col_gap-0.06,0.38,fs=8,c=NAVY,italic=True,align=PP_ALIGN.CENTER)
t(sl,"SPI/ADC",col_xs[1]+col_w3+0.04,arrow_y-0.2,col_gap-0.06,0.38,fs=8,c=TEAL,italic=True,align=PP_ALIGN.CENTER)

# ═══ SLIDE 27 (NEW) — UART Protocol Commands & Responses ═════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"UART ASCII Protocol — Commands & Responses","Complete reference for the 115200 baud 8N1 communication link")
txt(sl,[
    ("Protocol Structure",13,NAVY,True,False,None),
    ("All messages are ASCII text, terminated with CR+LF (\\r\\n). "
     "Commands are sent by the GUI; responses are sent by the STM32. "
     "Multiple responses may arrive per command (e.g. START produces many DATA lines).",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("DATA Message Format",13,GREEN,True,False,None),
    ("DATA,<freq_hz>,<magnitude_ohm>,<phase_deg>",11,BODY,False,False,None),
    ("Example:  DATA,5000,1234.56,-23.45",11,BODY,False,True,None),
    ("One DATA line per frequency step during a sweep.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("FLASH_DATA Message Format",13,ORANGE,True,False,None),
    ("FLASH_DATA,<freq_hz>,<|Z|x100>,<phase_x100>,<Re_x100>,<Im_x100>",11,BODY,False,False,None),
    ("All values are integers x100 (fixed-point). GUI divides by 100.0.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("STATUS Message Format",13,BLUE,True,False,None),
    ("STATUS,IDLE|RUNNING,START=n,STOP=n,STEP=n,RF=n.n,FLASH_COUNT=n",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("SWEEP Messages",13,TEAL,True,False,None),
    ("SWEEP,BEGIN — sent when sweep starts, before first DATA line.",11,BODY,False,False,None),
    ("SWEEP,DONE  — sent after last DATA line when sweep finishes.",11,BODY,False,False,None),
    ("",4,BODY,False,False,None),
    ("Error Messages",13,RED,True,False,None),
    ("ERROR,FLASH_FULL — flash region is full, data not written.",11,BODY,False,False,None),
    ("ERROR,FLASH_WRITE — flash write operation failed.",11,BODY,False,False,None),
],0.4,1.38,6.5,5.9)
# Right: message flow timeline
rbox(sl,7.0,1.45,6.0,5.8,WHITE,border=NAVY,bw=1)
t(sl,"Message Flow — Start Sweep",7.2,1.55,5.6,0.35,fs=11,c=NAVY,bold=True,align=PP_ALIGN.CENTER)
flow27=[("GUI sends:","START",BLUE),("STM32 replies:","SWEEP,BEGIN",TEAL),
        ("STM32 sends x99:","DATA,1000,4523.12,-2.34",GREEN),
        ("STM32 sends x99:","DATA,2000,3211.45,-5.67",GREEN),
        ("...","...",MUTED),
        ("STM32 replies:","SWEEP,DONE",TEAL),
        ("GUI sends:","STATUS",BLUE),
        ("STM32 replies:","STATUS,IDLE,START=1000,...",NAVY),
        ("GUI sends:","FLASH_STATUS",BLUE),
        ("STM32 replies:","FLASH_STATUS,COUNT=99,CAPACITY=390",ORANGE),
        ("GUI sends:","DUMP_FLASH",BLUE),
        ("STM32 sends x99:","FLASH_DATA,1000,452312,-234,...",ORANGE)]
fy=2.0; fh=0.38
for sender,msg,fill in flow27:
    lbox(sl,f"{sender}  {msg}",7.1,fy,5.8,fh,fill if fill!=MUTED else LIGHT,
         tc=WHITE if fill!=MUTED else MUTED,fs=8.5,bold=False,align=PP_ALIGN.LEFT)
    fy+=fh+0.03

# ═══ SLIDE 28 (NEW) — End-to-End Measurement Sequence ════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"End-to-End Measurement Sequence","Complete numbered flow from user click to updated Nyquist plot")
steps28=[
    ("User clicks Start Sweep",
     "GUI clears DataModel, clears plots, computes total expected points, sends 'START' via SerialWorker.",BLUE),
    ("STM32 receives START",
     "uart_comm.c calls Process_Start_Command(). FlashLog_BeginSweep() erases the 16 KB flash log region. "
     "Sends 'SWEEP,BEGIN'. Sets measurement_running=1.",NAVY),
    ("Sweep loop — for each frequency f",
     "AD9833_SetFrequency(f) writes freq word via SPI. HAL_Delay(100 ms) allows excitation to settle.",TEAL),
    ("ADC DMA acquisition",
     "HAL_ADC_Start_DMA() triggered. DMA1 Ch1 collects 512 half-words (CH1+CH2 interleaved, 256 each). "
     "DMA complete callback fires when done.",GREEN),
    ("Tone estimation",
     "impedance.c Estimate_Tone(ref_buf, f) and Estimate_Tone(sig_buf, f) solve the 2x2 LS system. "
     "Returns magnitude and phase for each channel.",PURPLE),
    ("BodePoint calculation",
     "|Z| = ref.magnitude / sig.magnitude * Rf. "
     "Phase = (sig.phaseRad - ref.phaseRad) * 180/pi. Re = |Z|*cos(phase). Im = |Z|*sin(phase).",ORANGE),
    ("Flash write + UART TX",
     "FlashLog_WritePoint(f, &result) stores the record to W25Q32 flash (page write, split across boundaries). "
     "UART sends 'DATA,f,|Z|,phase'.",RED),
    ("GUI receives DATA line",
     "_parse_data() splits CSV, calls DataModel.add_point(). PlotWidget.update_plots() redraws all three "
     "canvases. ConfigPanel shows current frequency and point count.",GREEN),
    ("Sweep completes",
     "After last frequency, STM32 sends 'SWEEP,DONE'. GUI status changes to IDLE. "
     "User can Export CSV or Dump Flash.",NAVY),
]
rw28=(W-0.8)/2; rh28=0.68
for i,(title,desc,fill) in enumerate(steps28):
    col=i%2; row=i//2
    sx=0.4+col*(rw28+0.1); sy=1.4+row*(rh28+0.1)
    lbox2(sl,f"Step {i+1}: {title}",desc,sx,sy,rw28,rh28,fill,tc=WHITE,tfs=10,sfs=8.5)

# ═══ SLIDE 29 — Q&A ═══════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Frequently Asked Questions","Key technical questions with concise answers")
qa=[("Why least-squares instead of FFT?",
     "The sweep step rarely aligns to a DFT bin; LS projection at the exact f eliminates spectral leakage and amplitude error."),
    ("What is the frequency accuracy?",
     "AD9833 resolution = 25 MHz/2^28 = 0.093 Hz. Accuracy tracks crystal tolerance (±20 ppm typ.)."),
    ("Why log data to flash during sweep?",
     "W25Q32 preserves results if USB disconnects mid-sweep; DUMP_FLASH replays them to the GUI afterward."),
    ("What is the ADC sample rate?",
     "72 MHz / ((61.5+12.5)x2) = 486 kHz total; 243 kHz per channel; alias-free BW ~121 kHz."),
    ("How is impedance magnitude computed?",
     "Z = (V_ref/V_sig) x Rf. Magnitude ratio of the two ADC tone estimates multiplied by feedback Rf."),
    ("Can it measure capacitive/inductive loads?",
     "Yes. Phase range is ±90°; Nyquist shows capacitive semicircles. Inductance gives positive Im(Z)."),
    ("What limits the upper frequency?",
     "Nyquist at 121 kHz. Practically 100 kHz for accuracy. AD9833 hardware can reach 12.5 MHz."),
    ("Is the GUI cross-platform?",
     "Yes — PyQt5, PySerial, and Matplotlib run on Windows, Linux, and macOS without code changes.")]
cw_q=(W-0.8)/2-0.1
for i,(q,a) in enumerate(qa):
    col=i%2; row=i//2
    qx=0.4+col*(cw_q+0.2); qy=1.4+row*1.45
    lbox(sl,f"Q{i+1}: {q}",qx,qy,cw_q,0.52,BLUE,tc=WHITE,fs=9.5,bold=True,align=PP_ALIGN.LEFT)
    lbox(sl,a,qx,qy+0.52,cw_q,0.83,CARD,tc=BODY,fs=9,bold=False,align=PP_ALIGN.LEFT)

# ═══ SLIDE 30 — Improvement Suggestions ══════════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,LIGHT)
hdr(sl,"Suggested Improvements","Potential enhancements for future hardware and software revisions")
suggestions=[
    ("Hardware",NAVY,[
        ("Anti-alias filter","Add ~120 kHz active LPF before each ADC input to suppress aliased noise."),
        ("Guarding electrode","Driven-guard ring eliminates stray capacitance in 3-electrode setups."),
        ("Galvanic isolation","Isolate USB ground from measurement ground to remove earth-loop interference."),
        ("Higher frequency range","Replace AD9833 with AD9959 or similar to extend sweep to 1 MHz+."),
        ("On-board calibration relay","Relay-switched precision R/C reference for automatic self-calibration."),
    ]),
    ("Firmware",TEAL,[
        ("DMA double-buffering","Overlap ADC acquisition with tone estimation to halve per-step sweep time."),
        ("Multi-frequency averaging","Average N sweeps per step to improve SNR in noisy environments."),
        ("FreeRTOS multitasking","Separate task for UART parsing keeps comms responsive during long sweeps."),
        ("Wireless upload","Integrate ESP32-C3 or BLE module for cable-free data export."),
        ("Equivalent circuit fitting","Implement Randles curve-fitting in Python GUI using scipy.optimize."),
    ]),
]
col_ww=(W-0.9)/2
for ci,(cat,col,items) in enumerate(suggestions):
    sx=0.4+ci*(col_ww+0.1); sy=1.42
    lbox(sl,cat,sx,sy,col_ww,0.5,col,tc=WHITE,fs=14,bold=True); iy=sy+0.52
    for j,(title,desc) in enumerate(items):
        lbox2(sl,f"{j+1}. {title}",desc,sx,iy,col_ww,0.82,LCARD,tc=BODY,tfs=10,sfs=9,border=col); iy+=0.84

# ═══ SLIDE 31 — Conclusion ════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLANK); bg(sl,NAVY)
oval(sl,9.8,4.5,5.0,5.0,TEAL); oval(sl,-0.8,-0.8,3.5,3.5,BLUE)
t(sl,"Conclusion",0.75,0.75,9.0,0.85,fs=34,c=WHITE,bold=True)
box(sl,0.75,1.58,7.0,0.04,TEAL)
txt(sl,[
    ("A complete, research-grade EIS instrument has been designed, implemented, and verified from silicon to software.",14,WHITE,False,False,None),
    ("",5,WHITE,False,False,None),
    ("The STM32F303CCT6 paired with the AD9833 DDS, OPA2140 TIA, and dual-channel ADC with DMA "
     "forms a compact portable signal chain that accurately characterises electrochemical impedance across 1 Hz to 100 kHz.",12,RGBColor(0xB3,0xC9,0xE8),False,False,None),
    ("",5,WHITE,False,False,None),
    ("The least-squares tone estimator replaces the naive DFT and delivers accurate complex impedance "
     "even when the sweep step does not fall on a Fourier bin — a critical correctness requirement.",12,RGBColor(0xB3,0xC9,0xE8),False,False,None),
    ("",5,WHITE,False,False,None),
    ("Non-volatile W25Q32 flash logging ensures no measurement data is lost. "
     "The Python desktop GUI provides real-time Bode and Nyquist visualisation, one-click CSV export, "
     "and full sweep control.",12,RGBColor(0xB3,0xC9,0xE8),False,False,None),
    ("",5,WHITE,False,False,None),
    ("The system is a solid foundation: anti-alias filtering, equivalent circuit fitting, "
     "and wireless export are natural next steps toward a commercially competitive EIS platform.",12,RGBColor(0xB3,0xC9,0xE8),False,False,None),
],0.75,1.68,8.9,5.6)
t(sl,"Key Achievements",10.1,1.65,2.9,0.42,fs=13,c=CYAN,bold=True)
ach_y=2.12
for a in ["Swept-sine EIS 1Hz-100kHz","Least-squares tone estimator","256-sample dual-ch ADC DMA",
          "Non-volatile flash logging","Live Bode + Nyquist GUI","USB-UART ASCII protocol","CSV export ready"]:
    t(sl,f"  {a}",10.1,ach_y,2.9,0.4,fs=11,c=RGBColor(0x4C,0xF6,0x9B)); ach_y+=0.44

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save("EIS_Project_Presentation.pptx")
print(f"Saved 31 slides -> EIS_Project_Presentation.pptx")
